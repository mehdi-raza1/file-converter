import streamlit as st
from PIL import Image
import io
import base64
from pathlib import Path
import pandas as pd
import json
import zipfile
from docx import Document
from docx.shared import Inches
import fitz  # PyMuPDF - Better PDF processing
from pdf2image import convert_from_bytes
import img2pdf
from pptx import Presentation
from pptx.util import Inches as PptxInches
import openpyxl
from reportlab.lib.pagesizes import letter, landscape
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Image as RLImage, Table, TableStyle, PageBreak
from reportlab.pdfgen import canvas as pdfcanvas
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.units import inch
from reportlab.lib import colors
from reportlab.lib.enums import TA_CENTER, TA_LEFT, TA_RIGHT, TA_JUSTIFY
from reportlab.lib.utils import ImageReader
from xml.sax.saxutils import escape
import logging
import traceback
import tempfile
import os
import sys

# Configure logging for production
logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(levelname)s - %(message)s',
    handlers=[
        logging.FileHandler('converter.log'),
        logging.StreamHandler(sys.stdout)
    ]
)
logger = logging.getLogger(__name__)

# Set page configuration
st.set_page_config(
    page_title="Universal File Converter Pro",
    page_icon="🔄",
    layout="wide",
    initial_sidebar_state="expanded"
)

# Production-ready error handling decorator
def handle_conversion_errors(func):
    """Decorator for handling conversion errors gracefully"""
    def wrapper(*args, **kwargs):
        try:
            logger.info(f"Starting conversion: {func.__name__}")
            result = func(*args, **kwargs)
            if result:
                logger.info(f"Conversion successful: {func.__name__}")
            else:
                logger.warning(f"Conversion returned None: {func.__name__}")
            return result
        except Exception as e:
            logger.error(f"Error in {func.__name__}: {str(e)}")
            logger.error(f"Traceback: {traceback.format_exc()}")
            st.error(f"Conversion failed: {str(e)}")
            return None
    return wrapper

# File size validation
def validate_file_size(uploaded_file, max_size_mb=50):
    """Validate file size before processing"""
    if uploaded_file is None:
        return False
    
    file_size = len(uploaded_file.getvalue())
    max_size_bytes = max_size_mb * 1024 * 1024
    
    if file_size > max_size_bytes:
        st.error(f"File size ({file_size / (1024*1024):.1f} MB) exceeds maximum allowed size ({max_size_mb} MB)")
        return False
    return True

# Memory management helper
def clear_memory():
    """Clear memory after heavy operations"""
    import gc
    gc.collect()

# Custom CSS with improved styling
st.markdown("""
<style>
    .main-header {
        text-align: center;
        padding: 2rem 0;
        background: linear-gradient(135deg, #667eea 0%, #764ba2 100%);
        color: white;
        border-radius: 10px;
        margin-bottom: 2rem;
    }
    .conversion-card {
        padding: 20px;
        border-radius: 10px;
        background: linear-gradient(135deg, #667eea 0%, #764ba2 100%);
        color: white;
        text-align: center;
        margin: 10px;
        box-shadow: 0 4px 6px rgba(0, 0, 0, 0.1);
    }
    .stButton>button {
        width: 100%;
        background-color: #4CAF50;
        color: white;
        padding: 12px;
        font-size: 16px;
        border-radius: 8px;
        border: none;
        transition: all 0.3s ease;
    }
    .stButton>button:hover {
        background-color: #45a049;
        transform: translateY(-2px);
        box-shadow: 0 4px 8px rgba(0, 0, 0, 0.2);
    }
    .success-message {
        padding: 1rem;
        background-color: #d4edda;
        border: 1px solid #c3e6cb;
        border-radius: 5px;
        color: #155724;
        margin: 1rem 0;
    }
    .error-message {
        padding: 1rem;
        background-color: #f8d7da;
        border: 1px solid #f5c6cb;
        border-radius: 5px;
        color: #721c24;
        margin: 1rem 0;
    }
    .info-box {
        padding: 1rem;
        background-color: #e3f2fd;
        border-left: 4px solid #2196f3;
        margin: 1rem 0;
        border-radius: 4px;
    }
</style>
""", unsafe_allow_html=True)

# Title with improved styling
st.markdown('<div class="main-header"><h1>🔄 Universal File Converter Pro</h1><p>Convert any file format with ease - Production Ready</p></div>', unsafe_allow_html=True)

# Conversion categories
conversion_categories = {
    "📄 To PDF": [
        "Word to PDF",
        "Excel to PDF", 
        "PowerPoint to PDF",
        "JPG to PDF",
        "PNG to PDF",
        "Text to PDF"
    ],
    "📝 From PDF": [
        "PDF to Word",
        "PDF to Excel",
        "PDF to PowerPoint",
        "PDF to JPG",
        "PDF to PNG",
        "Extract PDF Images",
        "PDF to Text"
    ],
    "🛠️ PDF Tools": [
        "Merge PDF",
        "Split PDF",
        "Compress PDF",
        "Rotate PDF",
        "Remove PDF Pages",
        "Extract PDF Pages"
    ],
    "🖼️ Image Conversion": [
        "JPG to PNG",
        "PNG to JPG",
        "Image to WebP",
        "WebP to JPG",
        "WebP to PNG",
        "Image to BMP",
        "BMP to JPG",
        "Resize Image",
        "Rotate Image"
    ],
    "📊 Office Files": [
        "Word to Excel",
        "Excel to Word",
        "CSV to Excel",
        "Excel to CSV",
        "JSON to Excel",
        "Excel to JSON"
    ]
}

# Sidebar
st.sidebar.header("🎯 Select Conversion Type")
category = st.sidebar.selectbox("Category", list(conversion_categories.keys()))
conversion_type = st.sidebar.selectbox("Conversion", conversion_categories[category])

# Helper Functions

@handle_conversion_errors
def word_to_pdf(uploaded_file):
    """Convert Word to PDF with enhanced error handling and formatting preservation"""
    if not validate_file_size(uploaded_file, 25):  # 25MB limit for Word files
        return None
    
    try:
        doc = Document(uploaded_file)
        
        output = io.BytesIO()
        pdf_doc = SimpleDocTemplate(output, pagesize=letter,
                                   leftMargin=inch, rightMargin=inch,
                                   topMargin=inch, bottomMargin=inch)
        
        styles = getSampleStyleSheet()
        story = []
        
        # Create enhanced custom styles
        for i in range(1, 10):
            style_name = f'Heading{i}'
            if style_name not in styles:
                styles.add(ParagraphStyle(
                    name=style_name,
                    parent=styles['Heading1'],
                    fontSize=max(12, 20 - (i * 2)),
                    spaceAfter=12,
                    spaceBefore=12,
                    textColor=colors.darkblue
                ))
        
        # Process paragraphs with enhanced error handling
        for para in doc.paragraphs:
            if not para.text.strip():
                story.append(Spacer(1, 6))
                continue
            
            style = styles['Normal']
            
            # Enhanced paragraph style detection with proper null checks
            if para.style and hasattr(para.style, 'name') and para.style.name and para.style.name.startswith('Heading'):
                try:
                    level = para.style.name.replace('Heading', '').strip()
                    if level.isdigit() and int(level) <= 9:
                        style = styles[f'Heading{level}']
                    else:
                        style = styles['Heading1']
                except (ValueError, KeyError, AttributeError):
                    style = styles['Heading1']
            elif para.style and hasattr(para.style, 'name') and para.style.name and 'Title' in para.style.name:
                style = styles['Title']
            
            # Escape special characters and limit text length
            text = escape(para.text[:2000])  # Limit text length
            if text.strip():
                p = Paragraph(text, style)
                story.append(p)
                story.append(Spacer(1, 6))
        
        # Process tables with enhanced formatting
        for table in doc.tables:
            try:
                data = []
                for row in table.rows:
                    row_data = []
                    for cell in row.cells:
                        cell_text = escape(cell.text[:500])  # Limit cell text
                        row_data.append(cell_text)
                    data.append(row_data)
                
                if data and any(any(cell.strip() for cell in row) for row in data):
                    # Calculate column widths
                    col_count = len(data[0]) if data else 1
                    available_width = pdf_doc.width * 0.9
                    col_widths = [available_width / col_count] * col_count
                    
                    t = Table(data, colWidths=col_widths, repeatRows=1)
                    t.setStyle(TableStyle([
                        ('BACKGROUND', (0, 0), (-1, 0), colors.lightblue),
                        ('TEXTCOLOR', (0, 0), (-1, 0), colors.black),
                        ('ALIGN', (0, 0), (-1, -1), 'LEFT'),
                        ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
                        ('FONTNAME', (0, 1), (-1, -1), 'Helvetica'),
                        ('FONTSIZE', (0, 0), (-1, -1), 10),
                        ('GRID', (0, 0), (-1, -1), 1, colors.black),
                        ('VALIGN', (0, 0), (-1, -1), 'TOP'),
                        ('LEFTPADDING', (0, 0), (-1, -1), 6),
                        ('RIGHTPADDING', (0, 0), (-1, -1), 6),
                    ]))
                    story.append(t)
                    story.append(Spacer(1, 12))
            except Exception as table_error:
                logger.warning(f"Table processing error: {table_error}")
                continue
        
        # Safety check for story length
        if len(story) > 1000:
            story = story[:1000]
            story.append(Paragraph("... (Content truncated for performance)", styles['Normal']))
        
        # Build PDF with fallback
        try:
            pdf_doc.build(story)
        except Exception as build_error:
            logger.warning(f"Complex layout failed: {build_error}")
            # Create simplified version
            simple_story = [
                Paragraph("Document Content", styles['Title']),
                Spacer(1, 20)
            ]
            
            for para in doc.paragraphs[:50]:  # Limit paragraphs
                if para.text.strip():
                    text = escape(para.text.strip()[:1000])
                    simple_story.append(Paragraph(text, styles['Normal']))
                    simple_story.append(Spacer(1, 8))
            
            pdf_doc.build(simple_story)
        
        output.seek(0)
        clear_memory()
        return output
        
    except Exception as e:
        logger.error(f"Word to PDF conversion error: {str(e)}")
        raise e

def excel_to_pdf(uploaded_file):
    """Convert Excel to PDF"""
    try:
        wb = openpyxl.load_workbook(uploaded_file)
        
        output = io.BytesIO()
        pdf_doc = SimpleDocTemplate(output, pagesize=letter)
        styles = getSampleStyleSheet()
        story = []
        
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            
            story.append(Paragraph(f"<b>Sheet: {sheet_name}</b>", styles['Heading1']))
            story.append(Spacer(1, 12))
            
            data = []
            for row in ws.iter_rows():
                row_data = [str(cell.value if cell.value is not None else "") for cell in row]
                data.append(row_data)
            
            if data:
                t = Table(data, repeatRows=1)
                t.setStyle(TableStyle([
                    ('BACKGROUND', (0, 0), (-1, 0), colors.HexColor('#4472C4')),
                    ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
                    ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
                    ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
                    ('GRID', (0, 0), (-1, -1), 0.5, colors.grey),
                ]))
                story.append(t)
                story.append(Spacer(1, 24))
        
        pdf_doc.build(story)
        output.seek(0)
        return output
    except Exception as e:
        st.error(f"Error: {str(e)}")
        return None

def ppt_to_pdf(uploaded_file):
    """Convert PowerPoint to PDF using a direct approach that preserves formatting"""
    try:
        # Validate file size
        file_content = uploaded_file.read()
        if len(file_content) > 50 * 1024 * 1024:  # 50MB limit
            raise ValueError("File size too large. Please upload a file smaller than 50MB.")
        
        # Reset file pointer
        uploaded_file.seek(0)
        
        # Create a temporary directory to work with the files
        import tempfile
        import os
        import subprocess
        import platform
        from pathlib import Path
        
        # Create a temporary directory
        with tempfile.TemporaryDirectory() as temp_dir:
            # Save the PowerPoint file to the temporary directory
            temp_ppt_path = os.path.join(temp_dir, "presentation.pptx")
            with open(temp_ppt_path, "wb") as f:
                f.write(file_content)
            
            # Create output PDF path
            temp_pdf_path = os.path.join(temp_dir, "output.pdf")
            
            # Try to use LibreOffice for conversion if available
            try:
                # Check if LibreOffice is installed
                if platform.system() == "Windows":
                    # Try common LibreOffice installation paths on Windows
                    libreoffice_paths = [
                        r"C:\Program Files\LibreOffice\program\soffice.exe",
                        r"C:\Program Files (x86)\LibreOffice\program\soffice.exe"
                    ]
                    
                    libreoffice_path = None
                    for path in libreoffice_paths:
                        if os.path.exists(path):
                            libreoffice_path = path
                            break
                    
                    if libreoffice_path:
                        # Use LibreOffice to convert PPT to PDF
                        cmd = [libreoffice_path, "--headless", "--convert-to", "pdf", "--outdir", temp_dir, temp_ppt_path]
                        subprocess.run(cmd, check=True, stdout=subprocess.PIPE, stderr=subprocess.PIPE)
                        
                        # Check if PDF was created
                        pdf_path = os.path.join(temp_dir, "presentation.pdf")
                        if os.path.exists(pdf_path):
                            with open(pdf_path, "rb") as pdf_file:
                                pdf_content = pdf_file.read()
                                return pdf_content
                
                # If LibreOffice conversion failed or not available, try comtypes approach on Windows
                if platform.system() == "Windows":
                    try:
                        import comtypes.client
                        
                        # Get absolute paths
                        abs_ppt_path = os.path.abspath(temp_ppt_path)
                        abs_pdf_path = os.path.abspath(temp_pdf_path)
                        
                        # Initialize PowerPoint application
                        powerpoint = comtypes.client.CreateObject("Powerpoint.Application")
                        powerpoint.Visible = False  # Run in background
                        
                        # Open the presentation
                        presentation = powerpoint.Presentations.Open(abs_ppt_path)
                        
                        # Save as PDF with high quality
                        # ppFixedFormatTypeXPS = 18, ppFixedFormatTypePDF = 17
                        # Use PDF format with high quality settings
                        fixedFormat = 17  # PDF format
                        
                        # Export with high quality settings
                        presentation.ExportAsFixedFormat(
                            abs_pdf_path,
                            fixedFormat,
                            PrintRange=1,  # All slides
                            OutputType=0,  # Standard (not handouts)
                            PrintHiddenSlides=False,
                            FrameSlides=False,
                            Intent=1,  # High quality
                            KeepIRMSettings=True
                        )
                        
                        # Close presentation and quit PowerPoint
                        presentation.Close()
                        powerpoint.Quit()
                        
                        # Force garbage collection to release COM objects
                        del presentation
                        del powerpoint
                        import gc
                        gc.collect()
                        
                        # Check if PDF was created
                        if os.path.exists(abs_pdf_path):
                            with open(abs_pdf_path, "rb") as pdf_file:
                                pdf_content = pdf_file.read()
                                return pdf_content
                        
                        raise Exception("PDF file not created by comtypes")
                    except Exception as comtypes_error:
                        st.warning(f"Windows COM conversion failed: {str(comtypes_error)}")
                
                # Try unoconv as another alternative
                try:
                    # Check if unoconv is installed
                    subprocess.run(["unoconv", "--version"], check=True, stdout=subprocess.PIPE, stderr=subprocess.PIPE)
                    
                    # Use unoconv to convert PPT to PDF
                    subprocess.run(["unoconv", "-f", "pdf", "-o", temp_dir, temp_ppt_path], check=True)
                    
                    # Check if PDF was created
                    pdf_path = os.path.join(temp_dir, "presentation.pdf")
                    if os.path.exists(pdf_path):
                        with open(pdf_path, "rb") as pdf_file:
                            pdf_content = pdf_file.read()
                            return pdf_content
                except:
                    # If all direct conversion methods failed, fall back to alternative method
                    pass
                
                raise Exception("All direct conversion methods failed")
                
            except Exception as e:
                # Fall back to python-pptx and reportlab approach with improved formatting
                st.warning("Using fallback conversion method. Some formatting may be affected.")
                
                # Load presentation
                prs = Presentation(temp_ppt_path)
                
                # Validate slide count
                if len(prs.slides) == 0:
                    raise ValueError("The PowerPoint file appears to be empty. Please check the file.")

                # High-fidelity fallback: absolute positioning with canvas to preserve layout
                try:
                    st.info("Using high-fidelity renderer (absolute positioning).")

                    # Utility: EMU to points
                    def emu_to_pt(val):
                        try:
                            return float(val) * 72.0 / 914400.0
                        except Exception:
                            return 0.0

                    page_width_pt = emu_to_pt(prs.slide_width)
                    page_height_pt = emu_to_pt(prs.slide_height)

                    output = io.BytesIO()
                    c = pdfcanvas.Canvas(output, pagesize=(page_width_pt, page_height_pt))

                    # Helper to draw background (try picture from slide/master, else solid color)
                    def draw_slide_background(slide):
                        # Try picture background via XML relationship
                        try:
                            from pptx.oxml.ns import qn
                            bg = slide._element.cSld.bg
                            if bg is not None and getattr(bg, 'bgPr', None) is not None:
                                bgPr = bg.bgPr
                                blipFill = getattr(bgPr, 'blipFill', None)
                                if blipFill is not None and getattr(blipFill, 'blip', None) is not None:
                                    blip = blipFill.blip
                                    rId = blip.get(qn('r:embed'))
                                    if rId:
                                        part = slide.part.related_parts.get(rId)
                                        if part is not None:
                                            pil_img = Image.open(io.BytesIO(part.blob))
                                            c.drawImage(ImageReader(pil_img), 0, 0, width=page_width_pt, height=page_height_pt)
                                            return
                        except Exception:
                            pass

                        # Try layout/master pictures (common in templates)
                        try:
                            # Slide Layout backgrounds/pictures
                            layout = getattr(slide, 'slide_layout', None)
                            if layout is not None:
                                for shp in getattr(layout, 'shapes', []):
                                    try:
                                        if hasattr(shp, 'image') and shp.image is not None:
                                            pil_img = Image.open(io.BytesIO(shp.image.blob))
                                            # Draw stretched as background (layout images often intended as full-bleed)
                                            c.drawImage(ImageReader(pil_img), 0, 0, width=page_width_pt, height=page_height_pt)
                                            return
                                    except Exception:
                                        continue
                            # Slide Master pictures
                            master = getattr(slide, 'slide_layout', None)
                            master = getattr(master, 'slide_master', None) if master else None
                            if master is not None:
                                for shp in getattr(master, 'shapes', []):
                                    try:
                                        if hasattr(shp, 'image') and shp.image is not None:
                                            pil_img = Image.open(io.BytesIO(shp.image.blob))
                                            c.drawImage(ImageReader(pil_img), 0, 0, width=page_width_pt, height=page_height_pt)
                                            return
                                    except Exception:
                                        continue
                        except Exception:
                            pass

                        # Fallback to solid fill color
                        try:
                            fill = slide.background.fill
                            if fill and hasattr(fill, 'foreground_color') and fill.fore_color and hasattr(fill.fore_color, 'rgb') and fill.fore_color.rgb:
                                rgb = fill.fore_color.rgb
                                r = rgb[0] / 255.0
                                g = rgb[1] / 255.0
                                b = rgb[2] / 255.0
                                c.setFillColor(colors.Color(r, g, b))
                                c.rect(0, 0, page_width_pt, page_height_pt, stroke=0, fill=1)
                            else:
                                # Default white
                                c.setFillColor(colors.white)
                                c.rect(0, 0, page_width_pt, page_height_pt, stroke=0, fill=1)
                        except Exception:
                            c.setFillColor(colors.white)
                            c.rect(0, 0, page_width_pt, page_height_pt, stroke=0, fill=1)

                    # Iterate slides and draw shapes with absolute positions
                    for slide in prs.slides:
                        draw_slide_background(slide)

                        for shape in slide.shapes:
                            try:
                                x_pt = emu_to_pt(getattr(shape, 'left', 0))
                                y_pt_top = emu_to_pt(getattr(shape, 'top', 0))
                                w_pt = emu_to_pt(getattr(shape, 'width', 0))
                                h_pt = emu_to_pt(getattr(shape, 'height', 0))
                                bottom_y = page_height_pt - y_pt_top - h_pt

                                # Draw pictures
                                if hasattr(shape, 'image') and shape.image is not None:
                                    try:
                                        img_blob = shape.image.blob
                                        img_stream = io.BytesIO(img_blob)
                                        pil_img = Image.open(img_stream)
                                        img_w, img_h = pil_img.size
                                        if img_w == 0 or img_h == 0:
                                            continue
                                        # Scale to fit bounding box preserving aspect ratio and center
                                        scale = min(w_pt / img_w, h_pt / img_h) if w_pt > 0 and h_pt > 0 else 1.0
                                        draw_w = img_w * scale
                                        draw_h = img_h * scale
                                        draw_x = x_pt + max(0, (w_pt - draw_w) / 2.0)
                                        draw_y = bottom_y + max(0, (h_pt - draw_h) / 2.0)
                                        c.drawImage(ImageReader(pil_img), draw_x, draw_y, width=draw_w, height=draw_h, preserveAspectRatio=True, mask='auto')
                                    except Exception:
                                        continue

                                # Draw text frames
                                elif getattr(shape, 'has_text_frame', False):
                                    try:
                                        tf = shape.text_frame
                                        from pptx.enum.text import PP_ALIGN
                                        # Vertical anchor
                                        try:
                                            from pptx.enum.text import MSO_ANCHOR
                                            v_anchor = getattr(tf, 'vertical_anchor', getattr(MSO_ANCHOR, 'TOP', None))
                                        except Exception:
                                            v_anchor = None

                                        # Build paragraphs preserving runs and formatting
                                        paras = []
                                        max_font_size = 12
                                        default_font_name = 'Helvetica'
                                        for p in tf.paragraphs:
                                            runs_html = []
                                            run_font_size = None
                                            run_font_name = None
                                            for r in getattr(p, 'runs', []) or []:
                                                txt = escape(getattr(r, 'text', '') or '')
                                                font = r.font
                                                size_pt = None
                                                name = None
                                                color_hex = None
                                                try:
                                                    if getattr(font, 'size', None):
                                                        size_pt = int(getattr(font.size, 'pt', 0)) or None
                                                        run_font_size = max(run_font_size or 0, size_pt or 0)
                                                        max_font_size = max(max_font_size, size_pt or max_font_size)
                                                    name = getattr(font, 'name', None)
                                                    if name and not run_font_name:
                                                        run_font_name = name
                                                    fc = getattr(font, 'color', None)
                                                    if fc is not None and getattr(fc, 'rgb', None):
                                                        rgb = fc.rgb
                                                        color_hex = '#%02x%02x%02x' % (rgb[0], rgb[1], rgb[2])
                                                except Exception:
                                                    pass
                                                open_tags = ''
                                                close_tags = ''
                                                if getattr(font, 'bold', False):
                                                    open_tags += '<b>'
                                                    close_tags = '</b>' + close_tags
                                                if getattr(font, 'italic', False):
                                                    open_tags += '<i>'
                                                    close_tags = '</i>' + close_tags
                                                if getattr(font, 'underline', False):
                                                    open_tags += '<u>'
                                                    close_tags = '</u>' + close_tags
                                                font_attrs = []
                                                if name:
                                                    font_attrs.append(f'name="{name}"')
                                                if size_pt:
                                                    font_attrs.append(f'size="{size_pt}"')
                                                if color_hex:
                                                    font_attrs.append(f'color="{color_hex}"')
                                                if font_attrs:
                                                    runs_html.append(f'<font {" ".join(font_attrs)}>{open_tags}{txt}{close_tags}</font>')
                                                else:
                                                    runs_html.append(f'{open_tags}{txt}{close_tags}')
                                            para_text = ''.join(runs_html) if runs_html else escape(getattr(p, 'text', '') or '')
                                            # Bullet support
                                            bullet_text = None
                                            level = getattr(p, 'level', 0) or 0
                                            if getattr(p, 'bullet', None) or level > 0:
                                                bullet_text = '•'
                                            # Alignment per paragraph
                                            align_map = {
                                                getattr(PP_ALIGN, 'LEFT', None): TA_LEFT,
                                                getattr(PP_ALIGN, 'CENTER', None): TA_CENTER,
                                                getattr(PP_ALIGN, 'RIGHT', None): TA_RIGHT,
                                                getattr(PP_ALIGN, 'JUSTIFY', None): TA_JUSTIFY,
                                            }
                                            p_align = align_map.get(getattr(p, 'alignment', None), TA_LEFT)
                                            # Style per paragraph
                                            left_indent = 18 * level  # points
                                            bullet_indent = max(0, left_indent - 12)
                                            # Line spacing from PPT if available
                                            try:
                                                line_spacing = getattr(p, 'line_spacing', None)
                                                leading = float(getattr(line_spacing, 'pt', 0)) if line_spacing else (run_font_size or max_font_size) * 1.2
                                            except Exception:
                                                leading = (run_font_size or max_font_size) * 1.2
                                            style = ParagraphStyle(
                                                name='ShapeText',
                                                fontName=(run_font_name or default_font_name),
                                                fontSize=(run_font_size or max_font_size),
                                                leading=leading,
                                                alignment=p_align,
                                                leftIndent=left_indent,
                                                bulletIndent=bullet_indent
                                            )
                                            paras.append((para_text, style, bullet_text))

                                        # Measure total height
                                        heights = []
                                        wrapped = []
                                        for txt, style, bullet in paras:
                                            para = Paragraph(txt, style, bulletText=bullet)
                                            wrap_w, wrap_h = para.wrap(w_pt, h_pt)
                                            heights.append(wrap_h)
                                            wrapped.append(para)
                                        total_h = sum(heights)
                                        # Determine vertical anchor
                                        anchor = getattr(tf, 'vertical_anchor', None)
                                        anchor_name = getattr(anchor, 'name', str(anchor)) if anchor is not None else ''
                                        if 'MIDDLE' in str(anchor_name):
                                            base_y = bottom_y + max(0, (h_pt - total_h)/2.0)
                                        elif 'BOTTOM' in str(anchor_name):
                                            base_y = bottom_y
                                        else:
                                            base_y = bottom_y + max(0, (h_pt - total_h))
                                        # Draw each paragraph stacked
                                        y_cursor = base_y
                                        for para, h in zip(wrapped, heights):
                                            para.drawOn(c, x_pt, y_cursor)
                                            y_cursor += h
                                    except Exception:
                                        continue

                                # Draw tables (approximate)
                                elif getattr(shape, 'has_table', False):
                                    try:
                                        data = []
                                        tbl = shape.table
                                        for row in tbl.rows:
                                            data.append([cell.text for cell in row.cells])
                                        # Calculate column widths evenly
                                        col_count = len(data[0]) if data else 0
                                        col_widths = [w_pt / max(col_count, 1)] * max(col_count, 1)
                                        t = Table(data, colWidths=col_widths)
                                        t.setStyle(TableStyle([
                                            ('GRID', (0,0), (-1,-1), 0.5, colors.grey),
                                            ('FONT', (0,0), (-1,-1), 'Helvetica', 12),
                                            ('VALIGN', (0,0), (-1,-1), 'MIDDLE')
                                        ]))
                                        tw, th = t.wrap(w_pt, h_pt)
                                        t.drawOn(c, x_pt, bottom_y + max(0, (h_pt - th)))
                                    except Exception:
                                        continue
                            except Exception:
                                # Skip problematic shapes
                                continue

                        c.showPage()

                    c.save()
                    output.seek(0)
                    return output
                except Exception as render_err:
                    # If high-fidelity rendering fails, continue with existing simple story builder below
                    st.warning(f"High-fidelity renderer fallback failed: {render_err}")
        
        # Get presentation dimensions for better page sizing
        try:
            # Try to get slide dimensions from presentation
            slide_width = prs.slide_width
            slide_height = prs.slide_height
            # Convert from EMU to inches (1 inch = 914400 EMU)
            width_inches = slide_width / 914400
            height_inches = slide_height / 914400
            
            # Create custom page size based on slide dimensions with some margins
            custom_pagesize = (width_inches * inch, height_inches * inch)
        except:
            # Fallback to letter size if dimensions can't be determined
            custom_pagesize = letter
        
        output = io.BytesIO()
        pdf_doc = SimpleDocTemplate(
            output,
            pagesize=custom_pagesize,
            leftMargin=0.5*inch,
            rightMargin=0.5*inch,
            topMargin=0.5*inch,
            bottomMargin=0.5*inch
        )
        styles = getSampleStyleSheet()
        story = []
        
        # Create enhanced custom styles with better formatting
        slide_title_style = ParagraphStyle(
            'SlideTitle',
            parent=styles['Heading1'],
            fontSize=20,
            spaceAfter=16,
            spaceBefore=8,
            alignment=TA_CENTER,
            textColor=colors.darkblue,
            fontName='Helvetica-Bold',
            leading=24  # Improved line spacing
        )
        
        content_title_style = ParagraphStyle(
            'ContentTitle',
            parent=styles['Heading2'],
            fontSize=16,
            spaceAfter=10,
            spaceBefore=8,
            alignment=TA_LEFT,
            textColor=colors.darkblue,
            fontName='Helvetica-Bold',
            leading=20  # Improved line spacing
        )
        
        content_style = ParagraphStyle(
            'SlideContent',
            parent=styles['Normal'],
            fontSize=12,
            spaceAfter=10,
            spaceBefore=4,
            alignment=TA_LEFT,
            leading=16,  # Improved line spacing
            fontName='Helvetica'
        )
        
        bullet_style = ParagraphStyle(
            'BulletStyle',
            parent=content_style,
            leftIndent=30,
            bulletIndent=15,
            spaceAfter=8,
            bulletFontName='Symbol',
            bulletText='•',
            leading=16  # Improved line spacing
        )
        
        # Try to extract presentation title from metadata or first slide
        presentation_title = "Photography Studios"  # Default to the filename
        try:
            if hasattr(prs, 'core_properties') and prs.core_properties.title:
                presentation_title = prs.core_properties.title
            elif prs.slides and hasattr(prs.slides[0], 'shapes'):
                for shape in prs.slides[0].shapes:
                    if hasattr(shape, 'text') and shape.text.strip() and len(shape.text) < 100:
                        presentation_title = shape.text.strip()
                        break
        except:
            pass
        
        # Add presentation title
        story.append(Paragraph(presentation_title, slide_title_style))
        story.append(Spacer(1, 20))
        
        # Limit number of slides to prevent excessive processing
        max_slides = 50  # Increased limit for better content coverage
        slides_to_process = list(prs.slides)[:max_slides]
        
        for i, slide in enumerate(slides_to_process):
            # Add slide header with better formatting
            slide_header = f"Slide {i + 1} of {len(slides_to_process)}"
            story.append(Paragraph(slide_header, slide_title_style))
            story.append(Spacer(1, 16))
            
            # Process slide content with increased limits
            shape_count = 0
            max_shapes_per_slide = 30  # Increased limit for better content capture
            slide_has_content = False
            
            # Try to extract slide title first
            slide_title = None
            try:
                # Look for title placeholder
                for shape in slide.shapes:
                    if hasattr(shape, 'is_placeholder') and shape.is_placeholder and shape.placeholder_format.idx == 0:
                        if hasattr(shape, 'text') and shape.text.strip():
                            slide_title = shape.text.strip()
                            break
                    # Fallback: look for any text that looks like a title
                    elif hasattr(shape, 'text') and shape.text.strip() and len(shape.text) < 100:
                        if shape.text.strip() not in slide_header:  # Avoid duplicating slide number
                            slide_title = shape.text.strip()
                            break
            except:
                pass
            
            # Add slide title if found
            if slide_title:
                story.append(Paragraph(slide_title, content_title_style))
                story.append(Spacer(1, 12))
                slide_has_content = True
            
            # Sort shapes by their position (top to bottom, left to right)
            try:
                sorted_shapes = sorted(slide.shapes, key=lambda s: (s.top, s.left) if hasattr(s, 'top') and hasattr(s, 'left') else (0, 0))
            except:
                sorted_shapes = slide.shapes
            
            for shape in sorted_shapes:
                # Skip if this is the title we already processed
                if slide_title and hasattr(shape, 'text') and shape.text.strip() == slide_title:
                    continue
                    
                shape_count += 1
                if shape_count > max_shapes_per_slide:
                    break
                
                try:
                    # Handle images with improved quality
                    if hasattr(shape, 'shape_type') and shape.shape_type == 13:  # Picture type
                        try:
                            # Extract image from shape
                            image = shape.image
                            image_bytes = image.blob
                            
                            # Create PIL Image
                            img_buffer = io.BytesIO(image_bytes)
                            pil_image = Image.open(img_buffer)
                            
                            # Calculate better image size based on PDF dimensions
                            max_width = min(pdf_doc.width * 0.8, 500)  # 80% of page width or 500pt max
                            max_height = min(pdf_doc.height * 0.6, 400)  # 60% of page height or 400pt max
                            
                            # Preserve aspect ratio
                            width, height = pil_image.size
                            aspect = width / height if height > 0 else 1
                            
                            if width > max_width:
                                width = max_width
                                height = width / aspect
                            
                            if height > max_height:
                                height = max_height
                                width = height * aspect
                            
                            # Ensure dimensions are valid
                            width = max(1, int(width))
                            height = max(1, int(height))
                            
                            # Resize image with high quality
                            try:
                                pil_image = pil_image.resize((width, height), Image.Resampling.LANCZOS)
                            except:
                                # Fallback to BICUBIC if LANCZOS fails
                                pil_image = pil_image.resize((width, height), Image.BICUBIC)
                            
                            # Convert to RGB if necessary
                            if pil_image.mode in ('RGBA', 'LA', 'P'):
                                pil_image = pil_image.convert('RGB')
                            
                            # Save to BytesIO with higher quality
                            img_buffer = io.BytesIO()
                            pil_image.save(img_buffer, format='JPEG', quality=95)
                            img_buffer.seek(0)
                            
                            # Create ReportLab image with proper alignment and preserve formatting
                            rl_image = RLImage(img_buffer, width=width, height=height)
                            # Center the image for better formatting
                            rl_image.hAlign = 'CENTER'
                            story.append(rl_image)
                            story.append(Spacer(1, 16))
                            slide_has_content = True
                            
                        except Exception as img_error:
                            # If image extraction fails, add a placeholder
                            story.append(Paragraph("[Image could not be extracted]", content_style))
                            story.append(Spacer(1, 8))
                            slide_has_content = True
                    
                    # Handle text content with improved formatting
                    elif hasattr(shape, 'text') and shape.text.strip():
                        text = shape.text.strip()
                        if len(text) > 2000:  # Increased text limit
                            text = text[:2000] + "..."
                        
                        # Escape special characters for ReportLab
                        text = text.replace('&', '&amp;').replace('<', '&lt;').replace('>', '&gt;')
                        
                        # Determine text type and apply appropriate styling
                        if len(text) < 100 and '\n' not in text:
                            # Likely a title or header
                            current_style = content_title_style
                        elif text.startswith(('•', '-', '*', '◦', '▪', '▫')):
                            # Bullet point
                            current_style = bullet_style
                            text = text[1:].strip()  # Remove bullet character
                        else:
                            current_style = content_style
                        
                        # Handle multi-line text better
                        paragraphs = text.split('\n')
                        for para in paragraphs:
                            if para.strip():
                                story.append(Paragraph(para.strip(), current_style))
                                story.append(Spacer(1, 4))
                        
                        story.append(Spacer(1, 8))
                        slide_has_content = True
                    
                    # Enhanced text frame processing with better formatting
                    elif hasattr(shape, 'text_frame') and shape.text_frame:
                        text_frame = shape.text_frame
                        if hasattr(text_frame, 'paragraphs'):
                            for paragraph in text_frame.paragraphs:
                                if paragraph.text.strip():
                                    para_text = paragraph.text.strip()
                                    if len(para_text) > 1500:
                                        para_text = para_text[:1500] + "..."
                                    
                                    # Escape special characters for ReportLab
                                    para_text = para_text.replace('&', '&amp;').replace('<', '&lt;').replace('>', '&gt;')
                                    
                                    # Enhanced bullet detection and formatting
                                    level = getattr(paragraph, 'level', 0)
                                    
                                    if (para_text.startswith(('•', '-', '*', '◦', '▪', '▫')) or level > 0):
                                        # Create custom bullet style based on level
                                        level_style = ParagraphStyle(
                                            f'BulletLevel{level}',
                                            parent=bullet_style,
                                            leftIndent=30 + (level * 15),
                                            bulletIndent=15 + (level * 15)
                                        )
                                        current_style = level_style
                                        
                                        if para_text.startswith(('•', '-', '*', '◦', '▪', '▫')):
                                            para_text = para_text[1:].strip()
                                    elif len(para_text) < 100 and '\n' not in para_text:
                                        current_style = content_title_style
                                    else:
                                        current_style = content_style
                                    
                                    story.append(Paragraph(para_text, current_style))
                                    story.append(Spacer(1, 4))
                                    slide_has_content = True
                    
                    # Enhanced table handling with better formatting
                    elif hasattr(shape, 'table'):
                        table = shape.table
                        table_data = []
                        
                        # Process table data with better formatting
                        for row_idx, row in enumerate(table.rows):
                            row_data = []
                            for cell in row.cells:
                                cell_text = cell.text.strip() if cell.text else ""
                                if len(cell_text) > 200:  # Increased cell text limit
                                    cell_text = cell_text[:200] + "..."
                                
                                # Escape special characters for ReportLab
                                cell_text = cell_text.replace('&', '&amp;').replace('<', '&lt;').replace('>', '&gt;')
                                row_data.append(cell_text)
                            
                            table_data.append(row_data)
                        
                        if table_data and any(any(cell for cell in row) for row in table_data):
                            # Calculate better column widths
                            col_count = len(table_data[0]) if table_data else 0
                            if col_count > 0:
                                available_width = pdf_doc.width * 0.9  # 90% of page width
                                col_widths = [available_width / col_count] * col_count
                                
                                # Create enhanced ReportLab table with better styling
                                t = Table(table_data, repeatRows=1, colWidths=col_widths)
                                t.setStyle(TableStyle([
                                    ('BACKGROUND', (0, 0), (-1, 0), colors.lightblue),
                                    ('TEXTCOLOR', (0, 0), (-1, 0), colors.black),
                                    ('ALIGN', (0, 0), (-1, -1), 'LEFT'),
                                    ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
                                    ('FONTNAME', (0, 1), (-1, -1), 'Helvetica'),
                                    ('FONTSIZE', (0, 0), (-1, -1), 10),  # Increased font size
                                    ('GRID', (0, 0), (-1, -1), 1, colors.black),
                                    ('VALIGN', (0, 0), (-1, -1), 'TOP'),
                                    ('ROWBACKGROUNDS', (0, 1), (-1, -1), [colors.white, colors.lightgrey]),
                                    ('LEFTPADDING', (0, 0), (-1, -1), 6),
                                    ('RIGHTPADDING', (0, 0), (-1, -1), 6),
                                    ('TOPPADDING', (0, 0), (-1, -1), 4),
                                    ('BOTTOMPADDING', (0, 0), (-1, -1), 4),
                                ]))
                                story.append(t)
                                story.append(Spacer(1, 16))
                                slide_has_content = True
                
                except Exception as shape_error:
                    # Skip problematic shapes but continue processing
                    continue
            
            # Add content if slide was empty
            if not slide_has_content:
                story.append(Paragraph("(Empty slide or content could not be extracted)", content_style))
                story.append(Spacer(1, 12))
            
            # Add separator between slides (except for last slide)
            if i < len(slides_to_process) - 1:
                story.append(Spacer(1, 20))
                story.append(PageBreak())
        
        # Safety check: limit total story elements
        if len(story) > 500:  # Increased limit for better content
            story = story[:500]
            story.append(Paragraph("... (Content truncated for safety - presentation too large)", content_style))
        
        # Build PDF with enhanced error handling
        try:
            pdf_doc.build(story)
        except Exception as build_error:
            # Enhanced fallback: create better formatted simple PDF
            st.warning("Complex layout failed, creating simplified PDF...")
            
            # Reset output buffer
            output.seek(0)
            output.truncate(0)
            
            # Use letter size for fallback
            pdf_doc = SimpleDocTemplate(
                output,
                pagesize=letter,
                leftMargin=0.75*inch,
                rightMargin=0.75*inch,
                topMargin=0.75*inch,
                bottomMargin=0.75*inch
            )
            
            simple_story = []
            simple_story.append(Paragraph(presentation_title, styles['Title']))
            simple_story.append(Spacer(1, 20))
            
            # Extract all content with better formatting
            for i, slide in enumerate(slides_to_process[:20]):  # Increased limit to 20 slides for fallback
                simple_story.append(Paragraph(f"Slide {i + 1}", styles['Heading2']))
                simple_story.append(Spacer(1, 12))
                
                # Try to extract slide title
                slide_title = None
                try:
                    for shape in slide.shapes:
                        if hasattr(shape, 'is_placeholder') and shape.is_placeholder and shape.placeholder_format.idx == 0:
                            if hasattr(shape, 'text') and shape.text.strip():
                                slide_title = shape.text.strip()
                                break
                except:
                    pass
                
                # Add slide title if found
                if slide_title:
                    simple_story.append(Paragraph(slide_title, styles['Heading3']))
                    simple_story.append(Spacer(1, 8))
                
                # Process shapes with better formatting
                for shape in slide.shapes:
                    try:
                        # Skip if this is the title we already processed
                        if slide_title and hasattr(shape, 'text') and shape.text.strip() == slide_title:
                            continue
                            
                        if hasattr(shape, 'text') and shape.text.strip():
                            text = shape.text.strip()[:1000]  # Increased text length
                            
                            # Escape special characters for ReportLab
                            text = text.replace('&', '&amp;').replace('<', '&lt;').replace('>', '&gt;')
                            
                            # Handle multi-line text better
                            paragraphs = text.split('\n')
                            for para in paragraphs:
                                if para.strip():
                                    simple_story.append(Paragraph(para.strip(), styles['Normal']))
                                    simple_story.append(Spacer(1, 4))
                            
                            simple_story.append(Spacer(1, 8))
                        elif hasattr(shape, 'shape_type') and shape.shape_type == 13:
                            simple_story.append(Paragraph("[Image present but not extracted in simplified mode]", styles['Normal']))
                            simple_story.append(Spacer(1, 6))
                    except:
                        continue
                
                simple_story.append(Spacer(1, 16))
                
                # Add page break between slides
                if i < len(slides_to_process[:20]) - 1:
                    simple_story.append(PageBreak())
            
            # Build simple version
            pdf_doc.build(simple_story)
        
        output.seek(0)
        return output
    except Exception as e:
        st.error(f"PowerPoint conversion error: {str(e)}")
        return None

def image_to_pdf(uploaded_file):
    """Convert Image to PDF"""
    try:
        img = Image.open(uploaded_file)
        
        if img.mode == 'RGBA':
            img = img.convert('RGB')
        
        output = io.BytesIO()
        img_byte_arr = io.BytesIO()
        img.save(img_byte_arr, format='JPEG')
        img_byte_arr.seek(0)
        
        pdf_bytes = img2pdf.convert(img_byte_arr.getvalue())
        output.write(pdf_bytes)
        output.seek(0)
        return output
    except Exception as e:
        st.error(f"Error: {str(e)}")
        return None

def text_to_pdf(text_content):
    """Convert Text to PDF"""
    try:
        output = io.BytesIO()
        pdf_doc = SimpleDocTemplate(output, pagesize=letter)
        styles = getSampleStyleSheet()
        story = []
        
        for line in text_content.split('\n'):
            if line.strip():
                story.append(Paragraph(line, styles['Normal']))
                story.append(Spacer(1, 12))
        
        pdf_doc.build(story)
        output.seek(0)
        return output
    except Exception as e:
        st.error(f"Error: {str(e)}")
        return None

@handle_conversion_errors
def pdf_to_word(uploaded_file):
    """Convert PDF to Word using PyMuPDF for better text extraction"""
    if not validate_file_size(uploaded_file):
        return None
    
    try:
        # Read PDF with PyMuPDF
        pdf_document = fitz.open(stream=uploaded_file.read(), filetype="pdf")
        doc = Document()
        
        # Add document title
        doc.add_heading('PDF Content', 0)
        
        for page_num in range(len(pdf_document)):
            page = pdf_document.load_page(page_num)
            
            # Extract text with better formatting
            text = page.get_text("text")
            
            if text.strip():
                # Add page header
                doc.add_heading(f'Page {page_num + 1}', level=1)
                
                # Split text into paragraphs
                paragraphs = text.split('\n\n')
                for para in paragraphs:
                    if para.strip():
                        doc.add_paragraph(para.strip())
            
            # Extract images from page
            try:
                image_list = page.get_images()
                for img_index, img in enumerate(image_list):
                    xref = img[0]
                    pix = fitz.Pixmap(pdf_document, xref)
                    
                    if pix.n - pix.alpha < 4:  # GRAY or RGB
                        img_data = pix.tobytes("png")
                        
                        # Save image to temporary file
                        with tempfile.NamedTemporaryFile(delete=False, suffix='.png') as tmp_file:
                            tmp_file.write(img_data)
                            tmp_file.flush()
                            
                            # Add image to document
                            try:
                                doc.add_picture(tmp_file.name, width=Inches(4))
                            except:
                                pass  # Skip if image can't be added
                            finally:
                                os.unlink(tmp_file.name)
                    
                    pix = None  # Release memory
            except Exception as img_error:
                logger.warning(f"Image extraction error on page {page_num}: {img_error}")
        
        pdf_document.close()
        
        output = io.BytesIO()
        doc.save(output)
        output.seek(0)
        clear_memory()
        return output
        
    except Exception as e:
        logger.error(f"PDF to Word conversion error: {str(e)}")
        raise e

@handle_conversion_errors
def pdf_to_excel(uploaded_file):
    """Convert PDF to Excel using PyMuPDF"""
    if not validate_file_size(uploaded_file):
        return None
    
    try:
        pdf_document = fitz.open(stream=uploaded_file.read(), filetype="pdf")
        all_data = []
        
        for page_num in range(len(pdf_document)):
            page = pdf_document.load_page(page_num)
            
            # Try to extract tables first
            try:
                tables = page.find_tables()
                if tables:
                    for table in tables:
                        table_data = table.extract()
                        for row in table_data:
                            all_data.append([f"Page {page_num + 1} - Table"] + list(row))
                        all_data.append([])  # Empty row between tables
                else:
                    # Extract text and split into rows
                    text = page.get_text("text")
                    lines = text.split('\n')
                    for line in lines:
                        if line.strip():
                            all_data.append([f"Page {page_num + 1}", line.strip()])
            except:
                # Fallback to text extraction
                text = page.get_text("text")
                lines = text.split('\n')
                for line in lines:
                    if line.strip():
                        all_data.append([f"Page {page_num + 1}", line.strip()])
        
        pdf_document.close()
        
        if not all_data:
            all_data = [["No content found"]]
        
        df = pd.DataFrame(all_data)
        output = io.BytesIO()
        df.to_excel(output, index=False, header=False, engine='openpyxl')
        output.seek(0)
        clear_memory()
        return output
        
    except Exception as e:
        logger.error(f"PDF to Excel conversion error: {str(e)}")
        raise e

@handle_conversion_errors
def pdf_to_text(uploaded_file):
    """Convert PDF to Text using PyMuPDF"""
    if not validate_file_size(uploaded_file):
        return None
    
    try:
        pdf_document = fitz.open(stream=uploaded_file.read(), filetype="pdf")
        text_content = ""
        
        for page_num in range(len(pdf_document)):
            page = pdf_document.load_page(page_num)
            text = page.get_text("text")
            
            if text.strip():
                text_content += f"\n--- Page {page_num + 1} ---\n"
                text_content += text + "\n\n"
        
        pdf_document.close()
        clear_memory()
        return text_content if text_content.strip() else "No text content found in PDF"
        
    except Exception as e:
        logger.error(f"PDF to Text conversion error: {str(e)}")
        raise e

def pdf_to_images(uploaded_file, format='JPEG'):
    """Convert PDF pages to images using PyMuPDF"""
    try:
        import fitz  # PyMuPDF
        
        pdf_document = fitz.open(stream=uploaded_file.read(), filetype="pdf")
        
        zip_buffer = io.BytesIO()
        with zipfile.ZipFile(zip_buffer, 'w', zipfile.ZIP_DEFLATED) as zip_file:
            for page_num in range(min(len(pdf_document), 20)):  # Limit to 20 pages
                page = pdf_document.load_page(page_num)
                
                # Convert page to image with good quality
                pix = page.get_pixmap(matrix=fitz.Matrix(2.0, 2.0))
                img_data = pix.tobytes("png")
                
                # Convert to PIL Image and then to desired format
                img = Image.open(io.BytesIO(img_data))
                if format == 'JPEG' and img.mode == 'RGBA':
                    img = img.convert('RGB')
                
                img_byte_arr = io.BytesIO()
                img.save(img_byte_arr, format=format)
                img_byte_arr.seek(0)
                zip_file.writestr(f'page_{page_num+1}.{format.lower()}', img_byte_arr.getvalue())
                
                pix = None  # Release memory
        
        pdf_document.close()
        zip_buffer.seek(0)
        return zip_buffer
    except Exception as e:
        st.error(f"Error: {str(e)}")
        return None

def merge_pdfs(uploaded_files):
    """Merge multiple PDFs using PyMuPDF"""
    try:
        import fitz  # PyMuPDF
        
        # Validate total file size
        total_size = sum(len(file.getvalue()) for file in uploaded_files)
        if total_size > 100 * 1024 * 1024:  # 100MB limit
            raise ValueError("Total file size too large. Please keep total size under 100MB.")
        
        merged_pdf = fitz.open()
        total_pages = 0
        
        for uploaded_file in uploaded_files:
            # Reset file pointer
            uploaded_file.seek(0)
            
            # Read PDF
            pdf_document = fitz.open(stream=uploaded_file.read(), filetype="pdf")
            current_pages = len(pdf_document)
            
            # Check page limit
            total_pages += current_pages
            if total_pages > 500:  # Limit total pages
                pdf_document.close()
                raise ValueError("Too many pages. Please keep total pages under 500.")
            
            # Insert all pages from this PDF
            merged_pdf.insert_pdf(pdf_document)
            pdf_document.close()
        
        # Write output
        output = io.BytesIO()
        merged_pdf.save(output)
        merged_pdf.close()
        output.seek(0)
        
        return output
    except Exception as e:
        st.error(f"Error: {str(e)}")
        return None

def split_pdf(uploaded_file, split_at):
    """Split PDF at specific page using PyMuPDF"""
    try:
        import fitz  # PyMuPDF
        
        pdf_document = fitz.open(stream=uploaded_file.read(), filetype="pdf")
        total_pages = len(pdf_document)
        
        zip_buffer = io.BytesIO()
        with zipfile.ZipFile(zip_buffer, 'w', zipfile.ZIP_DEFLATED) as zip_file:
            # Create first part (pages 0 to split_at-1)
            if split_at > 0:
                pdf_part1 = fitz.open()
                pdf_part1.insert_pdf(pdf_document, from_page=0, to_page=min(split_at-1, total_pages-1))
                
                output1 = io.BytesIO()
                pdf_part1.save(output1)
                pdf_part1.close()
                output1.seek(0)
                zip_file.writestr('part1.pdf', output1.getvalue())
            
            # Create second part (pages split_at to end)
            if split_at < total_pages:
                pdf_part2 = fitz.open()
                pdf_part2.insert_pdf(pdf_document, from_page=split_at, to_page=total_pages-1)
                
                output2 = io.BytesIO()
                pdf_part2.save(output2)
                pdf_part2.close()
                output2.seek(0)
                zip_file.writestr('part2.pdf', output2.getvalue())
        
        pdf_document.close()
        zip_buffer.seek(0)
        return zip_buffer
    except Exception as e:
        st.error(f"Error: {str(e)}")
        return None

@handle_conversion_errors
def compress_pdf(uploaded_file):
    """Compress PDF using PyMuPDF"""
    if not validate_file_size(uploaded_file):
        return None
    
    try:
        pdf_document = fitz.open(stream=uploaded_file.read(), filetype="pdf")
        
        # Compress by reducing image quality and removing unnecessary data
        for page_num in range(len(pdf_document)):
            page = pdf_document.load_page(page_num)
            
            # Get images and compress them
            image_list = page.get_images()
            for img_index, img in enumerate(image_list):
                xref = img[0]
                pix = fitz.Pixmap(pdf_document, xref)
                
                if pix.n - pix.alpha < 4:  # GRAY or RGB
                    # Compress image
                    img_data = pix.tobytes("jpeg", jpg_quality=70)
                    
                    # Replace image in PDF
                    pdf_document.update_stream(xref, img_data)
                
                pix = None
        
        # Save with compression
        output = io.BytesIO()
        pdf_document.save(output, garbage=4, deflate=True, clean=True)
        pdf_document.close()
        output.seek(0)
        clear_memory()
        return output
        
    except Exception as e:
        logger.error(f"PDF compression error: {str(e)}")
        raise e

@handle_conversion_errors
def rotate_pdf(uploaded_file, rotation):
    """Rotate PDF pages using PyMuPDF"""
    if not validate_file_size(uploaded_file):
        return None
    
    try:
        pdf_document = fitz.open(stream=uploaded_file.read(), filetype="pdf")
        
        for page_num in range(len(pdf_document)):
            page = pdf_document.load_page(page_num)
            page.set_rotation(rotation)
        
        output = io.BytesIO()
        pdf_document.save(output)
        pdf_document.close()
        output.seek(0)
        clear_memory()
        return output
        
    except Exception as e:
        logger.error(f"PDF rotation error: {str(e)}")
        raise e

@handle_conversion_errors
def remove_pdf_pages(uploaded_file, pages_to_remove):
    """Remove specific pages from PDF using PyMuPDF"""
    if not validate_file_size(uploaded_file):
        return None
    
    try:
        pdf_document = fitz.open(stream=uploaded_file.read(), filetype="pdf")
        
        # Convert to 0-based indexing and sort in reverse order
        pages_to_remove = sorted([int(p) - 1 for p in pages_to_remove], reverse=True)
        
        # Remove pages (in reverse order to maintain indices)
        for page_num in pages_to_remove:
            if 0 <= page_num < len(pdf_document):
                pdf_document.delete_page(page_num)
        
        if len(pdf_document) == 0:
            raise ValueError("Cannot remove all pages from PDF")
        
        output = io.BytesIO()
        pdf_document.save(output)
        pdf_document.close()
        output.seek(0)
        clear_memory()
        return output
        
    except Exception as e:
        logger.error(f"PDF page removal error: {str(e)}")
        raise e

@handle_conversion_errors
def extract_pdf_pages(uploaded_file, pages_to_extract):
    """Extract specific pages from PDF using PyMuPDF"""
    if not validate_file_size(uploaded_file):
        return None
    
    try:
        pdf_document = fitz.open(stream=uploaded_file.read(), filetype="pdf")
        new_pdf = fitz.open()
        
        # Convert to 0-based indexing
        pages_to_extract = [int(p) - 1 for p in pages_to_extract]
        
        # Extract pages
        for page_num in sorted(pages_to_extract):
            if 0 <= page_num < len(pdf_document):
                new_pdf.insert_pdf(pdf_document, from_page=page_num, to_page=page_num)
        
        if new_pdf.page_count == 0:
            raise ValueError("No valid pages found to extract")
        
        output = io.BytesIO()
        new_pdf.save(output)
        pdf_document.close()
        new_pdf.close()
        output.seek(0)
        clear_memory()
        return output
        
    except Exception as e:
        logger.error(f"PDF page extraction error: {str(e)}")
        raise e

def convert_image_format(uploaded_file, output_format):
    """Convert between image formats"""
    try:
        img = Image.open(uploaded_file)
        
        if output_format in ['JPEG', 'JPG'] and img.mode in ('RGBA', 'P'):
            img = img.convert('RGB')
        
        output = io.BytesIO()
        save_format = 'JPEG' if output_format == 'JPG' else output_format
        img.save(output, format=save_format)
        output.seek(0)
        return output
    except Exception as e:
        st.error(f"Error: {str(e)}")
        return None

def resize_image(uploaded_file, width, height, maintain_aspect=True):
    """Resize image"""
    try:
        img = Image.open(uploaded_file)
        
        if maintain_aspect:
            img.thumbnail((width, height), Image.Resampling.LANCZOS)
        else:
            img = img.resize((width, height), Image.Resampling.LANCZOS)
        
        output = io.BytesIO()
        img_format = img.format or 'PNG'
        if img_format == 'JPEG' and img.mode in ('RGBA', 'P'):
            img = img.convert('RGB')
        img.save(output, format=img_format)
        output.seek(0)
        return output
    except Exception as e:
        st.error(f"Error: {str(e)}")
        return None

def rotate_image(uploaded_file, angle):
    """Rotate image"""
    try:
        img = Image.open(uploaded_file)
        rotated = img.rotate(angle, expand=True)
        
        output = io.BytesIO()
        img_format = img.format or 'PNG'
        rotated.save(output, format=img_format)
        output.seek(0)
        return output
    except Exception as e:
        st.error(f"Error: {str(e)}")
        return None

def word_to_excel(uploaded_file):
    """Convert Word tables to Excel"""
    try:
        doc = Document(uploaded_file)
        
        all_data = []
        for table in doc.tables:
            for row in table.rows:
                row_data = [cell.text for cell in row.cells]
                all_data.append(row_data)
            all_data.append([])  # Empty row between tables
        
        if not all_data:
            st.warning("No tables found in document. Extracting text...")
            all_data = [[para.text] for para in doc.paragraphs if para.text.strip()]
        
        df = pd.DataFrame(all_data)
        output = io.BytesIO()
        df.to_excel(output, index=False, header=False, engine='openpyxl')
        output.seek(0)
        return output
    except Exception as e:
        st.error(f"Error: {str(e)}")
        return None

def excel_to_word(uploaded_file):
    """Convert Excel to Word"""
    try:
        df = pd.read_excel(uploaded_file)
        doc = Document()
        
        doc.add_heading('Excel Data', 0)
        
        table = doc.add_table(rows=len(df) + 1, cols=len(df.columns))
        table.style = 'Light Grid Accent 1'
        
        for i, column in enumerate(df.columns):
            table.rows[0].cells[i].text = str(column)
        
        for i, row in enumerate(df.itertuples(index=False), 1):
            for j, value in enumerate(row):
                table.rows[i].cells[j].text = str(value)
        
        output = io.BytesIO()
        doc.save(output)
        output.seek(0)
        return output
    except Exception as e:
        st.error(f"Error: {str(e)}")
        return None

def json_to_excel(json_data):
    """Convert JSON to Excel"""
    try:
        data = json.loads(json_data)
        
        if isinstance(data, list):
            df = pd.DataFrame(data)
        elif isinstance(data, dict):
            df = pd.DataFrame([data])
        else:
            st.error("Invalid JSON format")
            return None
        
        output = io.BytesIO()
        df.to_excel(output, index=False, engine='openpyxl')
        output.seek(0)
        return output
    except Exception as e:
        st.error(f"Error: {str(e)}")
        return None

def excel_to_json(uploaded_file):
    """Convert Excel to JSON"""
    try:
        df = pd.read_excel(uploaded_file)
        json_data = df.to_json(orient='records', indent=2)
        return json_data
    except Exception as e:
        st.error(f"Error: {str(e)}")
        return None

# Main conversion interface
st.markdown(f"### {conversion_type}")

# Conversion logic
if conversion_type == "Word to PDF":
    uploaded_file = st.file_uploader("Upload Word Document", type=['docx', 'doc'])
    if uploaded_file and st.button("Convert to PDF"):
        with st.spinner("Converting..."):
            result = word_to_pdf(uploaded_file)
            if result:
                st.success("✅ Conversion successful!")
                st.download_button("📥 Download PDF", result, f"{Path(uploaded_file.name).stem}.pdf", "application/pdf")

elif conversion_type == "Excel to PDF":
    uploaded_file = st.file_uploader("Upload Excel File", type=['xlsx', 'xls'])
    if uploaded_file and st.button("Convert to PDF"):
        with st.spinner("Converting..."):
            result = excel_to_pdf(uploaded_file)
            if result:
                st.success("✅ Conversion successful!")
                st.download_button("📥 Download PDF", result, f"{Path(uploaded_file.name).stem}.pdf", "application/pdf")

elif conversion_type == "PowerPoint to PDF":
    uploaded_file = st.file_uploader("Upload PowerPoint", type=['pptx', 'ppt'])
    if uploaded_file and st.button("Convert to PDF"):
        with st.spinner("Converting..."):
            result = ppt_to_pdf(uploaded_file)
            if result:
                st.success("✅ Conversion successful!")
                st.download_button("📥 Download PDF", result, f"{Path(uploaded_file.name).stem}.pdf", "application/pdf")

elif conversion_type in ["JPG to PDF", "PNG to PDF"]:
    file_type = conversion_type.split()[0].lower()
    uploaded_file = st.file_uploader(f"Upload {file_type.upper()} Image", type=[file_type, 'jpeg'])
    if uploaded_file and st.button("Convert to PDF"):
        with st.spinner("Converting..."):
            result = image_to_pdf(uploaded_file)
            if result:
                st.success("✅ Conversion successful!")
                st.download_button("📥 Download PDF", result, f"{Path(uploaded_file.name).stem}.pdf", "application/pdf")

elif conversion_type == "Text to PDF":
    text_input = st.text_area("Enter text to convert", height=300)
    if text_input and st.button("Convert to PDF"):
        with st.spinner("Converting..."):
            result = text_to_pdf(text_input)
            if result:
                st.success("✅ Conversion successful!")
                st.download_button("📥 Download PDF", result, "text_document.pdf", "application/pdf")

elif conversion_type == "PDF to Word":
    uploaded_file = st.file_uploader("Upload PDF", type=['pdf'])
    if uploaded_file and st.button("Convert to Word"):
        try:
            with st.spinner("Converting PDF to Word..."):
                result = pdf_to_word(uploaded_file)
            if result:
                st.success("✅ Conversion successful!")
                st.download_button("📥 Download Word", result, f"{Path(uploaded_file.name).stem}.docx", 
                                 "application/vnd.openxmlformats-officedocument.wordprocessingml.document")
        except Exception as e:
            st.error(f"❌ Conversion failed: {str(e)}")
            logger.error(f"PDF to Word conversion error: {str(e)}")

elif conversion_type == "PDF to Excel":
    uploaded_file = st.file_uploader("Upload PDF", type=['pdf'])
    if uploaded_file and st.button("Convert to Excel"):
        try:
            with st.spinner("Converting PDF to Excel..."):
                result = pdf_to_excel(uploaded_file)
            if result:
                st.success("✅ Conversion successful!")
                st.download_button("📥 Download Excel", result, f"{Path(uploaded_file.name).stem}.xlsx",
                                 "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet")
        except Exception as e:
            st.error(f"❌ Conversion failed: {str(e)}")
            logger.error(f"PDF to Excel conversion error: {str(e)}")

elif conversion_type == "PDF to PowerPoint":
    uploaded_file = st.file_uploader("Upload PDF", type=['pdf'])
    if uploaded_file and st.button("Convert to PowerPoint"):
        try:
            with st.spinner("Converting PDF to PowerPoint..."):
                result = pdf_to_ppt(uploaded_file)
            if result:
                st.success("✅ Conversion successful!")
                st.download_button("📥 Download PowerPoint", result, f"{Path(uploaded_file.name).stem}.pptx",
                                 "application/vnd.openxmlformats-officedocument.presentationml.presentation")
        except Exception as e:
            st.error(f"❌ Conversion failed: {str(e)}")
            logger.error(f"PDF to PowerPoint conversion error: {str(e)}")

elif conversion_type == "PDF to Text":
    uploaded_file = st.file_uploader("Upload PDF", type=['pdf'])
    if uploaded_file and st.button("Extract Text"):
        try:
            with st.spinner("Extracting text from PDF..."):
                result = pdf_to_text(uploaded_file)
            if result:
                st.success("✅ Text extracted successfully!")
                st.text_area("Extracted Text", result, height=300)
                st.download_button("📥 Download Text", result, f"{Path(uploaded_file.name).stem}.txt", "text/plain")
        except Exception as e:
            st.error(f"❌ Text extraction failed: {str(e)}")
            logger.error(f"PDF to Text conversion error: {str(e)}")

elif conversion_type in ["PDF to JPG", "PDF to PNG"]:
    format_type = conversion_type.split()[-1].upper()
    uploaded_file = st.file_uploader("Upload PDF", type=['pdf'])
    if uploaded_file and st.button(f"Convert to {format_type}"):
        try:
            with st.spinner(f"Converting PDF to {format_type} images..."):
                result = pdf_to_images(uploaded_file, format_type)
            if result:
                st.success("✅ Conversion successful! Multiple images will be downloaded as ZIP")
                st.download_button("📥 Download Images (ZIP)", result, f"{Path(uploaded_file.name).stem}_images.zip", "application/zip")
        except Exception as e:
            st.error(f"❌ Conversion failed: {str(e)}")
            logger.error(f"PDF to {format_type} conversion error: {str(e)}")

elif conversion_type == "Extract PDF Images":
    uploaded_file = st.file_uploader("Upload PDF", type=['pdf'])
    if uploaded_file and st.button("Extract Images"):
        try:
            with st.spinner("Extracting images from PDF..."):
                result = pdf_to_images(uploaded_file, 'PNG')
            if result:
                st.success("✅ Images extracted successfully!")
                st.download_button("📥 Download Images (ZIP)", result, f"{Path(uploaded_file.name).stem}_extracted.zip", "application/zip")
        except Exception as e:
            st.error(f"❌ Image extraction failed: {str(e)}")
            logger.error(f"PDF image extraction error: {str(e)}")

elif conversion_type == "Merge PDF":
    uploaded_files = st.file_uploader("Upload PDF files to merge", type=['pdf'], accept_multiple_files=True)
    if uploaded_files and len(uploaded_files) > 1 and st.button("Merge PDFs"):
        try:
            with st.spinner("Merging PDFs..."):
                result = merge_pdfs(uploaded_files)
            if result:
                st.success("✅ PDFs merged successfully!")
                st.download_button("📥 Download Merged PDF", result, "merged.pdf", "application/pdf")
        except Exception as e:
            st.error(f"❌ Merge failed: {str(e)}")
    elif uploaded_files and len(uploaded_files) == 1:
        st.warning("⚠️ Please upload at least 2 PDF files to merge")

elif conversion_type == "Split PDF":
    uploaded_file = st.file_uploader("Upload PDF", type=['pdf'])
    if uploaded_file:
        try:
            pdf_document = fitz.open(stream=uploaded_file.read(), filetype="pdf")
            total_pages = len(pdf_document)
            pdf_document.close()
            uploaded_file.seek(0)  # Reset file pointer
            st.info(f"📄 PDF has {total_pages} pages")
        except Exception as e:
            st.error(f"Error reading PDF: {str(e)}")
            total_pages = 0
        
        split_at = st.number_input("Split at page number", min_value=1, max_value=total_pages-1, value=1)
        
        if st.button("Split PDF"):
            with st.spinner("Splitting..."):
                result = split_pdf(uploaded_file, split_at)
                if result:
                    st.success("✅ PDF split successfully!")
                    st.download_button("📥 Download Split PDFs (ZIP)", result, f"{Path(uploaded_file.name).stem}_split.zip", "application/zip")

elif conversion_type == "Compress PDF":
    uploaded_file = st.file_uploader("Upload PDF", type=['pdf'])
    if uploaded_file and st.button("Compress PDF"):
        try:
            with st.spinner("Compressing PDF..."):
                result = compress_pdf(uploaded_file)
            if result:
                original_size = len(uploaded_file.getvalue())
                compressed_size = len(result.getvalue())
                reduction = ((original_size - compressed_size) / original_size) * 100
                
                st.success(f"✅ PDF compressed successfully!")
                st.info(f"Size reduction: {reduction:.1f}%")
                st.download_button("📥 Download Compressed PDF", result, f"{Path(uploaded_file.name).stem}_compressed.pdf", "application/pdf")
        except Exception as e:
            st.error(f"❌ Compression failed: {str(e)}")

elif conversion_type == "Rotate PDF":
    uploaded_file = st.file_uploader("Upload PDF", type=['pdf'])
    if uploaded_file:
        rotation = st.selectbox("Select rotation angle", [90, 180, 270])
        
        if st.button("Rotate PDF"):
            try:
                with st.spinner("Rotating PDF..."):
                    result = rotate_pdf(uploaded_file, rotation)
                if result:
                    st.success("✅ PDF rotated successfully!")
                    st.download_button("📥 Download Rotated PDF", result, f"{Path(uploaded_file.name).stem}_rotated.pdf", "application/pdf")
            except Exception as e:
                st.error(f"❌ Rotation failed: {str(e)}")

elif conversion_type == "Remove PDF Pages":
    uploaded_file = st.file_uploader("Upload PDF", type=['pdf'])
    if uploaded_file:
        try:
            pdf_document = fitz.open(stream=uploaded_file.read(), filetype="pdf")
            total_pages = len(pdf_document)
            pdf_document.close()
            uploaded_file.seek(0)  # Reset file pointer
            st.info(f"📄 PDF has {total_pages} pages")
            
            pages_input = st.text_input("Enter page numbers to remove (comma-separated, e.g., 1,3,5)")
            
            if st.button("Remove Pages"):
                try:
                    pages_to_remove = [int(p.strip()) for p in pages_input.split(',') if p.strip()]
                    with st.spinner("Removing pages..."):
                        result = remove_pdf_pages(uploaded_file, pages_to_remove)
                        if result:
                            st.success(f"✅ Removed {len(pages_to_remove)} pages successfully!")
                            st.download_button("📥 Download PDF", result, f"{Path(uploaded_file.name).stem}_modified.pdf", "application/pdf")
                except ValueError:
                    st.error("Please enter valid page numbers")
        except Exception as e:
            st.error(f"Error reading PDF: {str(e)}")

elif conversion_type == "Extract PDF Pages":
    uploaded_file = st.file_uploader("Upload PDF", type=['pdf'])
    if uploaded_file:
        try:
            pdf_document = fitz.open(stream=uploaded_file.read(), filetype="pdf")
            total_pages = len(pdf_document)
            pdf_document.close()
            uploaded_file.seek(0)  # Reset file pointer
            st.info(f"📄 PDF has {total_pages} pages")
            
            pages_input = st.text_input("Enter page numbers to extract (comma-separated, e.g., 1,3,5)")
            
            if st.button("Extract Pages"):
                try:
                    pages_to_extract = [int(p.strip()) for p in pages_input.split(',') if p.strip()]
                    with st.spinner("Extracting pages..."):
                        result = extract_pdf_pages(uploaded_file, pages_to_extract)
                        if result:
                            st.success(f"✅ Extracted {len(pages_to_extract)} pages successfully!")
                            st.download_button("📥 Download PDF", result, f"{Path(uploaded_file.name).stem}_extracted.pdf", "application/pdf")
                except ValueError:
                    st.error("Please enter valid page numbers")
        except Exception as e:
            st.error(f"Error reading PDF: {str(e)}")

elif conversion_type in ["JPG to PNG", "PNG to JPG"]:
    input_format = conversion_type.split()[0].lower()
    output_format = conversion_type.split()[-1].upper()
    
    uploaded_file = st.file_uploader(f"Upload {input_format.upper()} Image", type=[input_format, 'jpeg'])
    if uploaded_file and st.button(f"Convert to {output_format}"):
        with st.spinner("Converting..."):
            result = convert_image_format(uploaded_file, output_format)
            if result:
                st.success("✅ Conversion successful!")
                ext = output_format.lower()
                st.download_button("📥 Download Image", result, f"{Path(uploaded_file.name).stem}.{ext}", f"image/{ext}")

elif conversion_type == "Image to WebP":
    uploaded_file = st.file_uploader("Upload Image", type=['jpg', 'jpeg', 'png', 'bmp'])
    if uploaded_file and st.button("Convert to WebP"):
        with st.spinner("Converting..."):
            result = convert_image_format(uploaded_file, 'WEBP')
            if result:
                st.success("✅ Conversion successful!")
                st.download_button("📥 Download WebP", result, f"{Path(uploaded_file.name).stem}.webp", "image/webp")

elif conversion_type == "WebP to JPG":
    uploaded_file = st.file_uploader("Upload WebP Image", type=['webp'])
    if uploaded_file and st.button("Convert to JPG"):
        with st.spinner("Converting..."):
            result = convert_image_format(uploaded_file, 'JPEG')
            if result:
                st.success("✅ Conversion successful!")
                st.download_button("📥 Download JPG", result, f"{Path(uploaded_file.name).stem}.jpg", "image/jpeg")

elif conversion_type == "WebP to PNG":
    uploaded_file = st.file_uploader("Upload WebP Image", type=['webp'])
    if uploaded_file and st.button("Convert to PNG"):
        with st.spinner("Converting..."):
            result = convert_image_format(uploaded_file, 'PNG')
            if result:
                st.success("✅ Conversion successful!")
                st.download_button("📥 Download PNG", result, f"{Path(uploaded_file.name).stem}.png", "image/png")

elif conversion_type == "Image to BMP":
    uploaded_file = st.file_uploader("Upload Image", type=['jpg', 'jpeg', 'png', 'webp'])
    if uploaded_file and st.button("Convert to BMP"):
        with st.spinner("Converting..."):
            result = convert_image_format(uploaded_file, 'BMP')
            if result:
                st.success("✅ Conversion successful!")
                st.download_button("📥 Download BMP", result, f"{Path(uploaded_file.name).stem}.bmp", "image/bmp")

elif conversion_type == "BMP to JPG":
    uploaded_file = st.file_uploader("Upload BMP Image", type=['bmp'])
    if uploaded_file and st.button("Convert to JPG"):
        with st.spinner("Converting..."):
            result = convert_image_format(uploaded_file, 'JPEG')
            if result:
                st.success("✅ Conversion successful!")
                st.download_button("📥 Download JPG", result, f"{Path(uploaded_file.name).stem}.jpg", "image/jpeg")

elif conversion_type == "Resize Image":
    uploaded_file = st.file_uploader("Upload Image", type=['jpg', 'jpeg', 'png', 'webp', 'bmp'])
    if uploaded_file:
        img = Image.open(uploaded_file)
        st.image(img, caption=f"Original: {img.size[0]}x{img.size[1]} pixels", use_container_width=True)
        
        col1, col2 = st.columns(2)
        with col1:
            width = st.number_input("Width (pixels)", min_value=1, value=img.size[0])
        with col2:
            height = st.number_input("Height (pixels)", min_value=1, value=img.size[1])
        
        maintain_aspect = st.checkbox("Maintain aspect ratio", value=True)
        
        if st.button("Resize Image"):
            with st.spinner("Resizing..."):
                result = resize_image(uploaded_file, width, height, maintain_aspect)
                if result:
                    st.success("✅ Image resized successfully!")
                    resized_img = Image.open(result)
                    st.image(resized_img, caption=f"Resized: {resized_img.size[0]}x{resized_img.size[1]} pixels")
                    result.seek(0)
                    ext = Path(uploaded_file.name).suffix
                    st.download_button("📥 Download Resized Image", result, f"{Path(uploaded_file.name).stem}_resized{ext}", f"image/{ext[1:]}")

elif conversion_type == "Rotate Image":
    uploaded_file = st.file_uploader("Upload Image", type=['jpg', 'jpeg', 'png', 'webp', 'bmp'])
    if uploaded_file:
        img = Image.open(uploaded_file)
        st.image(img, caption="Original Image", use_container_width=True)
        
        angle = st.selectbox("Select rotation angle", [90, 180, 270, -90, -180, -270])
        
        if st.button("Rotate Image"):
            with st.spinner("Rotating..."):
                result = rotate_image(uploaded_file, angle)
                if result:
                    st.success("✅ Image rotated successfully!")
                    rotated_img = Image.open(result)
                    st.image(rotated_img, caption=f"Rotated {angle}°")
                    result.seek(0)
                    ext = Path(uploaded_file.name).suffix
                    st.download_button("📥 Download Rotated Image", result, f"{Path(uploaded_file.name).stem}_rotated{ext}", f"image/{ext[1:]}")

elif conversion_type == "Word to Excel":
    uploaded_file = st.file_uploader("Upload Word Document", type=['docx'])
    if uploaded_file and st.button("Convert to Excel"):
        with st.spinner("Converting..."):
            result = word_to_excel(uploaded_file)
            if result:
                st.success("✅ Conversion successful!")
                st.download_button("📥 Download Excel", result, f"{Path(uploaded_file.name).stem}.xlsx",
                                 "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet")

elif conversion_type == "Excel to Word":
    uploaded_file = st.file_uploader("Upload Excel File", type=['xlsx', 'xls'])
    if uploaded_file and st.button("Convert to Word"):
        with st.spinner("Converting..."):
            result = excel_to_word(uploaded_file)
            if result:
                st.success("✅ Conversion successful!")
                st.download_button("📥 Download Word", result, f"{Path(uploaded_file.name).stem}.docx",
                                 "application/vnd.openxmlformats-officedocument.wordprocessingml.document")

elif conversion_type == "CSV to Excel":
    uploaded_file = st.file_uploader("Upload CSV", type=['csv'])
    if uploaded_file and st.button("Convert to Excel"):
        with st.spinner("Converting..."):
            df = pd.read_csv(uploaded_file)
            output = io.BytesIO()
            df.to_excel(output, index=False, engine='openpyxl')
            output.seek(0)
            st.success("✅ Conversion successful!")
            st.download_button("📥 Download Excel", output, f"{Path(uploaded_file.name).stem}.xlsx",
                             "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet")

elif conversion_type == "Excel to CSV":
    uploaded_file = st.file_uploader("Upload Excel", type=['xlsx', 'xls'])
    if uploaded_file and st.button("Convert to CSV"):
        with st.spinner("Converting..."):
            df = pd.read_excel(uploaded_file)
            csv_data = df.to_csv(index=False)
            st.success("✅ Conversion successful!")
            st.download_button("📥 Download CSV", csv_data, f"{Path(uploaded_file.name).stem}.csv", "text/csv")

elif conversion_type == "JSON to Excel":
    json_input = st.text_area("Paste JSON data", height=300)
    
    st.markdown("**Or upload a JSON file:**")
    uploaded_file = st.file_uploader("Upload JSON file", type=['json'])
    
    if uploaded_file:
        json_input = uploaded_file.read().decode('utf-8')
        st.text_area("JSON Content", json_input, height=200)
    
    if json_input and st.button("Convert to Excel"):
        with st.spinner("Converting..."):
            result = json_to_excel(json_input)
            if result:
                st.success("✅ Conversion successful!")
                st.download_button("📥 Download Excel", result, "converted.xlsx",
                                 "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet")

elif conversion_type == "Excel to JSON":
    uploaded_file = st.file_uploader("Upload Excel", type=['xlsx', 'xls'])
    if uploaded_file and st.button("Convert to JSON"):
        with st.spinner("Converting..."):
            result = excel_to_json(uploaded_file)
            if result:
                st.success("✅ Conversion successful!")
                st.text_area("JSON Output", result, height=300)
                st.download_button("📥 Download JSON", result, f"{Path(uploaded_file.name).stem}.json", "application/json")

# Footer
st.markdown("---")
st.markdown("### 📚 Supported Conversions")
col1, col2, col3 = st.columns(3)

with col1:
    st.markdown("**To PDF:**")
    st.markdown("- Word, Excel, PowerPoint")
    st.markdown("- JPG, PNG, Text")

with col2:
    st.markdown("**From PDF:**")
    st.markdown("- Word, Excel, PowerPoint")
    st.markdown("- JPG, PNG, Text")
    st.markdown("- Extract Images")

with col3:
    st.markdown("**Other:**")
    st.markdown("- Merge/Split PDF")
    st.markdown("- Image Formats & Resize")
    st.markdown("- CSV ↔ Excel ↔ JSON")

st.markdown("---")
st.markdown("**💡 Tips:**")
st.markdown("- For best results with PDF conversions, use clear, high-quality source files")
st.markdown("- Image to PDF conversions maintain original image quality")
st.markdown("- PDF to Image requires poppler-utils to be installed")
st.markdown("- Large files may take longer to process")